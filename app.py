import streamlit as st
import requests
from bs4 import BeautifulSoup
from urllib.parse import urljoin, urlparse, urlunparse
from PIL import Image
import io
import hashlib
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import time
import re
from typing import List, Set, Tuple
import concurrent.futures
import zipfile

# Page config
st.set_page_config(page_title="Website Image Scraper", page_icon="🖼️", layout="wide")

# Constants
IMAGES_PER_PPT = 200  # 200 images per PowerPoint file
MAX_WORKERS = 5       # Concurrent downloads
TIMEOUT = 10
MAX_FILE_SIZE = 5 * 1024 * 1024  # 5MB max per image
IMAGES_PER_SLIDE = 4  # 4 images per slide (2x2 grid)

class UniversalImageScraper:
    def __init__(self, base_url: str):
        self.base_url = base_url
        self.domain = urlparse(base_url).netloc
        self.scheme = urlparse(base_url).scheme
        self.visited_urls: Set[str] = set()
        self.seen_hashes: Set[str] = set()
        self.all_images: List[str] = []
        self.last_debug_info: List[str] = []  # NEW: Store debug info
        self.session = requests.Session()
        self.session.headers.update({
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.7',
            'Accept-Language': 'en-US,en;q=0.9',
            'Accept-Encoding': 'gzip, deflate, br',
            'Referer': 'https://www.google.com/',
            'DNT': '1',
            'Connection': 'keep-alive',
            'Upgrade-Insecure-Requests': '1',
            'Sec-Fetch-Dest': 'document',
            'Sec-Fetch-Mode': 'navigate',
            'Sec-Fetch-Site': 'none',
            'Sec-Fetch-User': '?1',
            'Cache-Control': 'max-age=0',
        })
    
    def normalize_url(self, url: str) -> str:
        """Normalize URL to prevent duplicates"""
        parsed = urlparse(url)
        # Remove fragment and normalize
        return urlunparse((parsed.scheme, parsed.netloc, parsed.path, parsed.params, parsed.query, ''))
    
    def is_valid_url(self, url: str) -> bool:
        """Check if URL belongs to the same domain"""
        try:
            parsed = urlparse(url)
            return parsed.netloc == self.domain or parsed.netloc == f'www.{self.domain}' or parsed.netloc == self.domain.replace('www.', '')
        except:
            return False
    
    def extract_all_links(self, html: str, current_url: str) -> Set[str]:
        """Extract all links from HTML"""
        links = set()
        try:
            soup = BeautifulSoup(html, 'html.parser')
            for tag in soup.find_all('a', href=True):
                href = tag['href']
                full_url = urljoin(current_url, href)
                normalized = self.normalize_url(full_url)
                
                if self.is_valid_url(normalized) and normalized not in self.visited_urls:
                    links.add(normalized)
        except:
            pass
        return links
    
    def extract_all_images(self, html: str, page_url: str) -> List[str]:
        """Extract ALL images from HTML - no filtering"""
        images = []
        debug_info = []
        
        try:
            soup = BeautifulSoup(html, 'html.parser')
            
            # CRITICAL: Check for Shopify JSON product data first
            script_count = 0
            for script in soup.find_all('script'):
                script_count += 1
                script_content = script.string or ''
                
                if 'cdn.shopify.com' in script_content:
                    debug_info.append(f"Found Shopify CDN in script {script_count}")
                    # Extract ALL cdn.shopify.com URLs
                    shopify_urls = re.findall(r'(https?:)?//cdn\.shopify\.com[^\s"\'<>]+', script_content, re.IGNORECASE)
                    debug_info.append(f"Extracted {len(shopify_urls)} Shopify URLs from script")
                    
                    for url in shopify_urls:
                        if url.startswith('//'):
                            url = 'https:' + url
                        # Clean the URL
                        url = url.split('"')[0].split("'")[0].split(',')[0].split(')')[0]
                        images.append(url)
            
            debug_info.append(f"Total scripts found: {script_count}")
            debug_info.append(f"Images from scripts: {len(images)}")
            
            # Method 1: <img> tags with all possible attributes
            img_count = 0
            for img in soup.find_all('img'):
                img_count += 1
                # Try multiple attributes
                src = (img.get('src') or 
                       img.get('data-src') or 
                       img.get('data-lazy-src') or 
                       img.get('data-original') or
                       img.get('data-lazy') or
                       img.get('data-srcset') or
                       img.get('data-fallback-src'))
                
                if src:
                    full_url = urljoin(page_url, src)
                    images.append(full_url)
                
                # Handle srcset attribute
                srcset = img.get('srcset')
                if srcset:
                    for src_entry in srcset.split(','):
                        url = src_entry.strip().split()[0]
                        full_url = urljoin(page_url, url)
                        images.append(full_url)
            
            debug_info.append(f"<img> tags found: {img_count}")
            
            # Method 2: Extract ALL URLs from entire HTML and filter for images
            all_urls = re.findall(r'https?://[^\s"\'<>]+', html)
            all_urls.extend(re.findall(r'//cdn\.[^\s"\'<>]+', html))
            
            debug_info.append(f"Total URLs in HTML: {len(all_urls)}")
            
            url_images = 0
            for url in all_urls:
                url = url.split('"')[0].split("'")[0].split(',')[0].split(')')[0]
                if url.startswith('//'):
                    url = 'https:' + url
                if self.is_image_url(url):
                    images.append(url)
                    url_images += 1
            
            debug_info.append(f"Image URLs from regex: {url_images}")
            
        except Exception as e:
            debug_info.append(f"ERROR: {str(e)}")
        
        # Filter to keep only valid image URLs and remove duplicates
        valid_images = []
        seen = set()
        for img_url in images:
            if img_url and img_url not in seen:
                if self.is_image_url(img_url):
                    valid_images.append(img_url)
                    seen.add(img_url)
        
        # Store debug info
        self.last_debug_info = debug_info
        
        return valid_images
    
    def is_image_url(self, url: str) -> bool:
        """Check if URL points to an image - very permissive"""
        if not url or len(url) < 4:
            return False
        
        try:
            url_lower = url.lower()
            
            # CRITICAL: Shopify CDN detection
            if 'cdn.shopify.com' in url_lower:
                return True
            
            # Remove query parameters for extension check
            url_without_query = url_lower.split('?')[0].split('#')[0]
            
            # Check for common image extensions
            image_extensions = [
                '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', 
                '.svg', '.ico', '.tiff', '.tif', '.avif', '.jfif',
                '.pjpeg', '.pjp', '.apng', '.heic', '.heif'
            ]
            
            if any(url_without_query.endswith(ext) for ext in image_extensions):
                return True
            
            # Check for image-related keywords in path
            image_keywords = [
                '/image/', '/img/', '/photo/', '/picture/', '/pic/', 
                '/media/', '/asset/', '/upload/', '/content/', 
                '/gallery/', '/thumbnail/', '/thumb/', '/banner/',
                '/icon/', '/logo/', '/bg/', '/background/', '/files/'
            ]
            
            if any(keyword in url_lower for keyword in image_keywords):
                return True
            
            # Check for image CDN patterns
            cdn_patterns = [
                'cloudinary', 'imgix', 'cloudflare', 'fastly',
                'akamai', 'images', 'static', 'cdn', 'assets'
            ]
            
            if any(pattern in url_lower for pattern in cdn_patterns):
                return True
            
            # If URL has image in query params (e.g., ?image=...)
            if 'image' in url_lower or 'img' in url_lower or 'photo' in url_lower:
                return True
                
        except:
            pass
        
        return False
    
    def extract_images_aggressive(self, html: str, page_url: str) -> List[str]:
        """Aggressive extraction - finds any URL that looks like an image"""
        images = []
        
        try:
            # Method 1: Extract from Shopify JSON data
            # Shopify stores product data in <script> tags with type="application/json"
            soup = BeautifulSoup(html, 'html.parser')
            
            # Look for Shopify product JSON
            for script in soup.find_all('script', type='application/json'):
                try:
                    script_content = script.string
                    if script_content:
                        # Find all URLs in the JSON that look like images
                        json_urls = re.findall(r'https?://[^"\s]+\.(?:jpg|jpeg|png|gif|webp|svg)[^"\s]*', script_content, re.IGNORECASE)
                        json_urls.extend(re.findall(r'//cdn\.shopify\.com/[^"\s]+', script_content))
                        
                        for url in json_urls:
                            if url.startswith('//'):
                                url = self.scheme + ':' + url
                            images.append(url.strip('",\'();[]{}'))
                except:
                    pass
            
            # Method 2: Find all URLs in the HTML using regex
            all_urls = re.findall(r'https?://[^\s<>"\']+', html)
            all_urls.extend(re.findall(r'//[^\s<>"\']+', html))  # Protocol-relative URLs
            all_urls.extend(re.findall(r'/[^\s<>"\']+\.(jpg|jpeg|png|gif|webp|svg|ico)', html, re.IGNORECASE))
            
            for url in all_urls:
                # Clean up the URL
                url = url.strip('",\'();[]{}')
                
                # Convert protocol-relative URLs
                if url.startswith('//'):
                    url = self.scheme + ':' + url
                
                # Make absolute
                full_url = urljoin(page_url, url)
                
                # Check if it's an image
                if self.is_image_url(full_url):
                    images.append(full_url)
            
            # Method 3: Look specifically for Shopify CDN URLs
            shopify_pattern = r'(https?:)?//cdn\.shopify\.com/[^"\s<>\']+\.(jpg|jpeg|png|gif|webp)'
            shopify_urls = re.findall(shopify_pattern, html, re.IGNORECASE)
            for match in shopify_urls:
                url = match[0] + '//cdn.shopify.com/' if not match[0] else match[0] + '//cdn.shopify.com/'
                # Reconstruct the full match
                full_match = ''.join([x for x in match if x])
                if full_match.startswith('//'):
                    full_match = self.scheme + ':' + full_match
                images.append(full_match)
        
        except Exception as e:
            pass
        
        # Remove duplicates and return
        return list(set([img for img in images if img]))

    
    def crawl_website(self, status_container, progress_callback) -> List[str]:
        """Recursively crawl entire website and collect all image URLs"""
        to_visit = {self.base_url}
        total_images = []
        page_count = 0
        
        status_container.write(f"🌐 Starting deep crawl of {self.domain}...")
        
        while to_visit and page_count < 1000:  # Safety limit
            current_batch = list(to_visit)[:10]  # Process 10 pages at a time
            to_visit -= set(current_batch)
            
            for url in current_batch:
                if url in self.visited_urls:
                    continue
                
                try:
                    self.visited_urls.add(url)
                    page_count += 1
                    
                    status_container.write(f"📄 Crawling page {page_count}: {url[:80]}...")
                    
                    response = self.session.get(url, timeout=TIMEOUT, allow_redirects=True)
                    response.raise_for_status()
                    
                    html = response.text
                    
                    # Extract images
                    page_images = self.extract_all_images(html, url)
                    
                    # Show debug info
                    if self.last_debug_info:
                        for debug_msg in self.last_debug_info:
                            status_container.write(f"      🔍 DEBUG: {debug_msg}")
                    
                    # Fallback: if no images found, try aggressive extraction
                    if not page_images:
                        status_container.write(f"   🔍 No images with standard methods, trying aggressive extraction...")
                        page_images = self.extract_images_aggressive(html, url)
                    
                    total_images.extend(page_images)
                    
                    status_container.write(f"   ✅ Found {len(page_images)} images on this page (Total: {len(total_images)})")
                    
                    # Extract new links
                    new_links = self.extract_all_links(html, url)
                    to_visit.update(new_links)
                    
                    progress_callback(page_count, len(total_images), len(to_visit))
                    
                    time.sleep(0.1)  # Be polite
                    
                except Exception as e:
                    status_container.write(f"   ⚠️ Error: {str(e)[:60]}")
        
        # Remove duplicates and filter
        unique_images = list(set(total_images))
        status_container.write(f"✅ Crawl complete! Found {len(unique_images)} unique image URLs across {page_count} pages")
        return unique_images
    
    def download_and_validate_image(self, url: str) -> Tuple[str, bytes] or None:
        """Download and validate a single image"""
        try:
            response = self.session.get(url, timeout=TIMEOUT, stream=True)
            response.raise_for_status()
            
            # Check content type
            content_type = response.headers.get('content-type', '').lower()
            if 'image' not in content_type and not self.is_image_url(url):
                return None
            
            # Check size
            content_length = response.headers.get('content-length')
            if content_length and int(content_length) > MAX_FILE_SIZE:
                return None
            
            # Download
            img_bytes = b''
            for chunk in response.iter_content(chunk_size=8192):
                img_bytes += chunk
                if len(img_bytes) > MAX_FILE_SIZE:
                    return None
            
            # Check hash
            img_hash = hashlib.sha256(img_bytes).hexdigest()
            if img_hash in self.seen_hashes:
                return None
            
            # Validate as image
            try:
                img = Image.open(io.BytesIO(img_bytes))
                
                # Skip tiny images
                if img.width < 50 or img.height < 50:
                    return None
                
                # Resize if too large
                max_dimension = 1920
                if img.width > max_dimension or img.height > max_dimension:
                    ratio = min(max_dimension/img.width, max_dimension/img.height)
                    new_size = (int(img.width * ratio), int(img.height * ratio))
                    img = img.resize(new_size, Image.Resampling.LANCZOS)
                
                # Convert to RGB
                if img.mode in ('RGBA', 'LA', 'P'):
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode == 'P':
                        img = img.convert('RGBA')
                    if img.mode in ('RGBA', 'LA'):
                        background.paste(img, mask=img.split()[-1])
                    else:
                        background.paste(img)
                    img = background
                elif img.mode != 'RGB':
                    img = img.convert('RGB')
                
                # Save to bytes
                output = io.BytesIO()
                img.save(output, format='JPEG', quality=85, optimize=True)
                final_bytes = output.getvalue()
                
                self.seen_hashes.add(img_hash)
                return (url, final_bytes)
                
            except:
                return None
                
        except:
            return None
    
    def download_all_images(self, image_urls: List[str], status_container) -> List[Tuple[str, bytes]]:
        """Download all images with progress tracking"""
        valid_images = []
        total = len(image_urls)
        
        status_container.write(f"⬇️ Starting download of {total} images...")
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
            futures = {executor.submit(self.download_and_validate_image, url): url for url in image_urls}
            
            completed = 0
            for future in concurrent.futures.as_completed(futures):
                completed += 1
                result = future.result()
                
                if result:
                    valid_images.append(result)
                    if len(valid_images) % 20 == 0:
                        status_container.write(f"   ✅ Downloaded {len(valid_images)} valid images ({completed}/{total} processed)")
        
        status_container.write(f"✅ Successfully downloaded {len(valid_images)} valid images!")
        return valid_images

def generate_ppt(images: List[Tuple[str, bytes]], batch_num: int, total_batches: int, domain: str) -> bytes:
    """Generate a single PowerPoint file"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank_layout = prs.slide_layouts[6]
    
    # Process images in groups of 4
    for i in range(0, len(images), IMAGES_PER_SLIDE):
        batch = images[i:i + IMAGES_PER_SLIDE]
        slide = prs.slides.add_slide(blank_layout)
        
        # White background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Grid layout
        num_images = len(batch)
        if num_images == 1:
            grid = [(0, 0)]
            cols, rows = 1, 1
        elif num_images == 2:
            grid = [(0, 0), (1, 0)]
            cols, rows = 2, 1
        elif num_images == 3:
            grid = [(0, 0), (1, 0), (0, 1)]
            cols, rows = 2, 2
        else:
            grid = [(0, 0), (1, 0), (0, 1), (1, 1)]
            cols, rows = 2, 2
        
        slide_width = prs.slide_width.inches
        slide_height = prs.slide_height.inches
        margin = 0.3
        h_spacing = 0.2
        v_spacing = 0.2
        
        available_width = slide_width - (2 * margin) - ((cols - 1) * h_spacing)
        available_height = slide_height - (2 * margin) - ((rows - 1) * v_spacing)
        
        cell_width = available_width / cols
        cell_height = available_height / rows
        
        for idx, (url, img_bytes) in enumerate(batch):
            col, row = grid[idx]
            
            img = Image.open(io.BytesIO(img_bytes))
            img_width, img_height = img.size
            img_ratio = img_width / img_height
            
            if img_ratio > (cell_width / cell_height):
                width = Inches(cell_width)
                height = Inches(cell_width / img_ratio)
            else:
                height = Inches(cell_height)
                width = Inches(cell_height * img_ratio)
            
            left = Inches(margin + (col * (cell_width + h_spacing)) + (cell_width - width.inches) / 2)
            top = Inches(margin + (row * (cell_height + v_spacing)) + (cell_height - height.inches) / 2)
            
            img_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(img_stream, left, top, width, height)
    
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()

def create_all_ppts(images: List[Tuple[str, bytes]], domain: str, status_container) -> bytes:
    """Create multiple PPT files and package in ZIP"""
    num_batches = (len(images) + IMAGES_PER_PPT - 1) // IMAGES_PER_PPT
    
    if num_batches == 1:
        # Single PPT
        status_container.write(f"📊 Generating PowerPoint with {len(images)} images...")
        ppt_bytes = generate_ppt(images, 1, 1, domain)
        status_container.write(f"✅ PowerPoint generated!")
        return ppt_bytes
    
    # Multiple PPTs in ZIP
    status_container.write(f"📦 Creating {num_batches} PowerPoint files...")
    zip_buffer = io.BytesIO()
    
    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
        for batch_idx in range(num_batches):
            start_idx = batch_idx * IMAGES_PER_PPT
            end_idx = min(start_idx + IMAGES_PER_PPT, len(images))
            batch_images = images[start_idx:end_idx]
            
            status_container.write(f"📊 Generating PPT {batch_idx + 1}/{num_batches} (images {start_idx + 1}-{end_idx})...")
            
            ppt_bytes = generate_ppt(batch_images, batch_idx + 1, num_batches, domain)
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{domain.replace('.', '_')}_batch_{batch_idx + 1}_of_{num_batches}_{timestamp}.pptx"
            
            zip_file.writestr(filename, ppt_bytes)
            
            num_slides = (len(batch_images) + IMAGES_PER_SLIDE - 1) // IMAGES_PER_SLIDE
            status_container.write(f"✅ Completed PPT {batch_idx + 1}/{num_batches} ({len(batch_images)} images, {num_slides} slides)")
    
    status_container.write(f"🎉 All {num_batches} PowerPoint files created!")
    zip_buffer.seek(0)
    return zip_buffer.getvalue()

def main():
    st.title("🖼️ Universal Website Image Scraper")
    st.markdown("**Download ALL images from any website** • 200 images per PowerPoint • Automatic batching")
    
    with st.expander("ℹ️ How It Works", expanded=False):
        st.info("""
        **Complete Website Scraping:**
        - Crawls the ENTIRE website recursively
        - Extracts ALL images (JPG, PNG, GIF, WebP, SVG, etc.)
        - No filtering - gets every single image
        - Downloads and validates all images
        - Creates PPT files with 200 images each (4 per slide)
        - Automatically packages multiple PPTs in ZIP
        
        **Example:** 
        - Website has 857 images → Creates 5 PPT files:
          - Batch 1: 200 images (50 slides)
          - Batch 2: 200 images (50 slides)
          - Batch 3: 200 images (50 slides)
          - Batch 4: 200 images (50 slides)
          - Batch 5: 57 images (15 slides)
        
        **Note:** Large websites may take several minutes to crawl.
        """)
    
    with st.expander("🧪 Test URLs", expanded=False):
        st.markdown("""
        **Try these example sites:**
        - `https://example.com` (Small, simple site)
        - `https://unsplash.com` (Photo gallery)
        - Any public website URL
        
        **Tips:**
        - Use the homepage URL (e.g., `https://example.com`)
        - Include `www.` if the site uses it
        - Make sure the site is publicly accessible
        - Be patient - large sites take longer
        """)
    
    url = st.text_input("🌐 Enter Website URL:", placeholder="https://example.com")
    
    if st.button("🚀 Scrape ALL Images", type="primary"):
        if not url:
            st.error("Please enter a valid URL")
            return
        
        if not url.startswith(('http://', 'https://')):
            url = 'https://' + url
        
        progress_bar = st.progress(0)
        status_container = st.status("Starting complete website scrape...", expanded=True)
        
        try:
            scraper = UniversalImageScraper(url)
            
            # Progress callback
            def update_progress(pages, images, remaining):
                progress = min(0.5 * (pages / max(pages + remaining, 1)), 0.5)
                progress_bar.progress(progress)
            
            # Step 1: Crawl website and collect image URLs
            status_container.write("🔍 Phase 1: Crawling website and discovering images...")
            status_container.write(f"Target domain: {scraper.domain}")
            status_container.write(f"Starting URL: {url}")
            progress_bar.progress(0.1)
            
            image_urls = scraper.crawl_website(status_container, update_progress)
            
            if not image_urls:
                status_container.update(label="❌ No images found", state="error")
                st.error(f"""
                No images found on this website. Possible reasons:
                - The website may block automated scraping
                - The URL might be incorrect
                - The site might not have any images
                - Try a different page on the site
                
                Debug info:
                - Pages crawled: {len(scraper.visited_urls)}
                - Image URLs found: 0
                """)
                return
            
            status_container.write(f"📊 Phase 1 Complete: Found {len(image_urls)} image URLs")
            
            # Step 2: Download and validate all images
            status_container.write(f"⬇️ Phase 2: Downloading {len(image_urls)} images...")
            progress_bar.progress(0.5)
            
            valid_images = scraper.download_all_images(image_urls, status_container)
            
            if not valid_images:
                status_container.update(label="❌ No valid images downloaded", state="error")
                st.error("All images failed validation or download.")
                return
            
            status_container.write(f"✅ Phase 2 Complete: {len(valid_images)} images ready")
            
            # Step 3: Generate PPTs
            status_container.write("📊 Phase 3: Generating PowerPoint presentations...")
            progress_bar.progress(0.8)
            
            result_bytes = create_all_ppts(valid_images, scraper.domain, status_container)
            
            progress_bar.progress(1.0)
            status_container.update(label="✅ Complete!", state="complete")
            
            # Display results
            num_ppts = (len(valid_images) + IMAGES_PER_PPT - 1) // IMAGES_PER_PPT
            
            if num_ppts == 1:
                st.success(f"🎉 Created PowerPoint with {len(valid_images)} images!")
                
                col1, col2 = st.columns(2)
                with col1:
                    st.metric("Total Images", len(valid_images))
                with col2:
                    st.metric("Total Slides", (len(valid_images) + IMAGES_PER_SLIDE - 1) // IMAGES_PER_SLIDE)
                
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = f"{scraper.domain.replace('.', '_')}_all_images_{timestamp}.pptx"
                
                st.download_button(
                    label="📥 Download PowerPoint",
                    data=result_bytes,
                    file_name=filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    type="primary"
                )
            else:
                st.success(f"🎉 Created {num_ppts} PowerPoint files with {len(valid_images)} total images!")
                
                col1, col2, col3 = st.columns(3)
                with col1:
                    st.metric("Total Images", len(valid_images))
                with col2:
                    st.metric("PPT Files", num_ppts)
                with col3:
                    st.metric("ZIP Size", f"{len(result_bytes) / 1024 / 1024:.1f} MB")
                
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = f"{scraper.domain.replace('.', '_')}_all_images_{num_ppts}_files_{timestamp}.zip"
                
                st.download_button(
                    label=f"📥 Download ZIP ({num_ppts} PPT files)",
                    data=result_bytes,
                    file_name=filename,
                    mime="application/zip",
                    type="primary"
                )
                
                with st.expander("📋 What's in the ZIP?"):
                    for i in range(num_ppts):
                        start = i * IMAGES_PER_PPT + 1
                        end = min((i + 1) * IMAGES_PER_PPT, len(valid_images))
                        slides = (end - start + 1 + IMAGES_PER_SLIDE - 1) // IMAGES_PER_SLIDE
                        st.write(f"- `batch_{i+1}.pptx` - Images {start}-{end} ({slides} slides)")
            
            # Preview
            st.subheader("🖼️ Preview (First 8 Images)")
            cols = st.columns(8)
            for i, (url, img_bytes) in enumerate(valid_images[:8]):
                with cols[i]:
                    st.image(img_bytes, use_container_width=True)
        
        except Exception as e:
            status_container.update(label=f"❌ Error: {str(e)}", state="error")
            st.error(f"Fatal error: {str(e)}")
            st.exception(e)

if __name__ == "__main__":
    main()
